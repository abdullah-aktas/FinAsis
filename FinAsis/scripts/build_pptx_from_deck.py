"""
Build a PPTX from the Markdown investor deck.

Parses docs/FinAsis_Investor_Deck_TR_v2.md and generates
docs/slides/FinAsis_Investor_Deck_TR_v2.pptx with one slide per top-level section (## ...).

Rules:
- Lines starting with '## ' start a new slide (Title + Content layout).
- Lines starting with '### ' become bold bullets.
- List items ('- ', '* ') become bullets (level 0).
- Blockquotes ('> ') and normal non-empty lines become subtle bullets.
- Code fences and mermaid blocks are skipped.
"""
from __future__ import annotations

import os
import re
from pathlib import Path

from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT
from typing import Optional, Any

ROOT = Path(__file__).resolve().parents[1]
DECK_MD = ROOT / "docs" / "FinAsis_Investor_Deck_TR_v2.md"
OUT_DIR = ROOT / "docs" / "slides"
OUT_PPTX = OUT_DIR / "FinAsis_Investor_Deck_TR_v2.pptx"


def sanitize(text: str) -> str:
    # Strip markdown inline artifacts
    text = text.strip()
    text = re.sub(r"`+", "", text)
    text = re.sub(r"\*\*(.*?)\*\*", r"\1", text)
    text = re.sub(r"\*(.*?)\*", r"\1", text)
    text = re.sub(r"\[(.*?)\]\((.*?)\)", r"\1", text)
    return text.strip()


def add_bullet(shape: Any, text: str, level: int = 0, bold: bool = False):
    tf: Any = getattr(shape, "text_frame", None)
    if tf is None:
        return
    p = tf.add_paragraph()
    p.text = text
    p.level = max(0, min(level, 5))
    p.font.size = Pt(18)
    p.font.bold = bold
    p.alignment = PP_PARAGRAPH_ALIGNMENT.LEFT


def set_shape_text(shape: Any, text: str) -> None:
    if shape is None:
        return
    tf: Any = getattr(shape, "text_frame", None)
    if tf is not None:
        try:
            tf.clear()  # pyright: ignore[reportAttributeAccessIssue]
        except Exception:
            pass
        try:
            tf.text = text  # pyright: ignore[reportAttributeAccessIssue]
        except Exception:
            try:
                if getattr(tf, "paragraphs", None):
                    tf.paragraphs[0].text = text  # pyright: ignore[reportAttributeAccessIssue]
            except Exception:
                pass
        return
    if hasattr(shape, "text"):
        try:
            shape.text = text  # pyright: ignore[reportAttributeAccessIssue]
        except Exception:
            pass


def add_footer_box(slide, text: str):
    try:
        width = slide.part.presentation.slide_width
        height = slide.part.presentation.slide_height
        # bottom area
        left = width - Inches(2.2)
        top = height - Inches(0.5)
        box_w = Inches(2.1)
        box_h = Inches(0.35)
        tx = slide.shapes.add_textbox(left, top, box_w, box_h)
        tf = tx.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(12)
        p.alignment = PP_PARAGRAPH_ALIGNMENT.RIGHT
    except Exception:
        # non-fatal
        pass


def main():
    if not DECK_MD.exists():
        raise FileNotFoundError(f"Deck markdown not found: {DECK_MD}")
    OUT_DIR.mkdir(parents=True, exist_ok=True)

    # Optional corporate template support
    template_path = OUT_DIR / "template.pptx"
    prs = Presentation(str(template_path)) if template_path.exists() else Presentation()
    title_layout = prs.slide_layouts[0]
    title_content_layout = prs.slide_layouts[1]

    # Cover slide
    slide = prs.slides.add_slide(title_layout)
    title_shape = slide.shapes.title
    set_shape_text(title_shape, "FinAsis – Yatırımcı Sunumu v2 (TR)")
    # Alt başlık placeholder'ı opsiyonel olabilir
    subtitle_shape: Any = slide.placeholders[1] if len(slide.placeholders) > 1 else None
    set_shape_text(subtitle_shape, "Görünürlük + Öngörü + Otomasyon\nGeleceğin zenginliğini inşa eden finans katmanı")

    in_code = False
    current_shape: Optional[Any] = None

    with DECK_MD.open("r", encoding="utf-8") as f:
        for raw in f:
            line = raw.rstrip("\n")

            # code fences skip
            if line.strip().startswith("```"):
                in_code = not in_code
                continue
            if in_code:
                continue

            if line.startswith("## "):
                title = sanitize(line[3:])
                slide = prs.slides.add_slide(title_content_layout)
                sh_title = slide.shapes.title
                set_shape_text(sh_title, title)
                # İçerik placeholder'ı genelde 1. index, tip güvenli kontrolle alalım
                current_shape = slide.shapes.placeholders[1] if len(slide.placeholders) > 1 else None
                tf_any: Any = getattr(current_shape, "text_frame", None)
                if tf_any is not None:
                    try:
                        tf_any.clear()  # pyright: ignore[reportAttributeAccessIssue]
                    except Exception:
                        pass
                # start with empty paragraph removed
                continue

            if current_shape is None:
                # ignore until first section
                continue

            if not line.strip():
                # blank line
                continue

            # Subsection as bold bullet
            if line.startswith("### "):
                add_bullet(current_shape, sanitize(line[4:]), level=0, bold=True)
                continue

            # List items (unordered or numbered) with nested levels based on indentation
            stripped = line.lstrip()
            indent = len(line) - len(stripped)
            list_level = max(0, min(5, indent // 2))  # 2 spaces per level (tolerant)
            if stripped.startswith("- ") or stripped.startswith("* "):
                content = sanitize(stripped[2:])
                add_bullet(current_shape, content, level=list_level, bold=False)
                continue
            m = re.match(r"(\d+)\.\s+(.*)", stripped)
            if m:
                content = sanitize(m.group(2))
                add_bullet(current_shape, content, level=list_level, bold=False)
                continue

            # Blockquote → subtle bullet
            if stripped.startswith("> "):
                # Blockquote level based on number of '>' and indentation
                m_gt = re.match(r"^>+", stripped)
                gt_count = len(m_gt.group(0)) if m_gt else 1
                content = sanitize(stripped.lstrip("> "))
                add_bullet(current_shape, content, level=min(5, 1 + (gt_count - 1) + (indent // 2)), bold=False)
                continue

            # Otherwise: treat as a small bullet
            add_bullet(current_shape, sanitize(line), level=1, bold=False)

    # Add slide numbers and footer to all but cover
    total = len(prs.slides)
    for idx, s in enumerate(prs.slides, start=1):
        footer = f"FinAsis  •  {idx}/{total}" if idx > 1 else "FinAsis"
        add_footer_box(s, footer)

    prs.save(str(OUT_PPTX))
    print(f"✅ PPTX created: {OUT_PPTX}")


if __name__ == "__main__":
    main()
